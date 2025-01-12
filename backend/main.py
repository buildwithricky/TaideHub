from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import json
import logging
import traceback
from dotenv import load_dotenv
from fastapi.responses import FileResponse
import google.generativeai as genai
from typing import Dict
import asyncio
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('app.log')
    ]
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Configure Google Gemini
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY not found in environment variables")

genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-pro')

# Initialize FastAPI app
app = FastAPI()

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000", "http://127.0.0.1:5173", "http://127.0.0.1:3000"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

async def generate_presentation_content(topic: str) -> dict:
    """Generate presentation content using Google Gemini."""
    logger.info(f"Generating presentation content for topic: {topic}")
    
    prompt = f"""Create an educational presentation about {topic} suitable for secondary school students.
    Return ONLY a JSON array with exactly 5 slides, following this exact format:
    [
        {{
            "title": "Title of slide 1",
            "subtitle": "Optional subtitle",
            "content": "• First bullet point\\n• Second bullet point"
        }},
        {{
            "title": "Let's Get Started!",
            "subtitle": "Starter Activity",
            "content": "• Think-Pair-Share Activity:\\n• Question: [engaging question about {topic}]\\n• Take 2 minutes to think\\n• Discuss with your partner\\n• Share with the class"
        }},
        {{
            "title": "Main Content",
            "subtitle": "Understanding [topic aspect]",
            "content": "• Detailed point 1 with explanation\\n• Detailed point 2 with explanation\\n• Knowledge Check: [specific question about the content]"
        }},
        {{
            "title": "Exploring Further",
            "subtitle": "Real-world Applications",
            "content": "• Detailed example 1 with explanation\\n• Detailed example 2 with explanation\\n• Knowledge Check: [application-based question]"
        }},
        {{
            "title": "Plenary",
            "subtitle": "Let's Review",
            "content": "• What have we learned today?\\n• Key takeaways:\\n  - [Key point 1]\\n  - [Key point 2]\\n• Exit Ticket: [final check for understanding]"
        }}
    ]
    
    Make sure to:
    1. Include exactly 5 slides
    2. Make content suitable for secondary school students
    3. Include engaging starter activities
    4. Add knowledge check questions
    5. Use detailed explanations
    6. End with a plenary review
    7. Use valid JSON format with double quotes
    8. Use \\n between bullet points
    9. Start bullet points with •"""

    try:
        logger.info("Making API call to Google Gemini")
        response = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    temperature=0.7,  # Increased for more creative responses
                    candidate_count=1,
                    stop_sequences=["],"],
                    top_p=0.8,
                    top_k=40,
                )
            )
        )
        
        # Extract and clean the response
        content_str = response.text.strip()
        logger.info(f"Raw response from Gemini: {content_str}")
        
        # Ensure the response is wrapped in square brackets
        if not content_str.startswith('['):
            content_str = '[' + content_str
        if not content_str.endswith(']'):
            content_str = content_str + ']'
            
        # Clean up any markdown formatting that might be present
        content_str = content_str.replace('```json', '').replace('```', '')
        
        try:
            content = json.loads(content_str)
            logger.info("Successfully parsed JSON content")
            return content
        except json.JSONDecodeError as je:
            logger.error(f"JSON parsing error: {str(je)}")
            logger.error(f"Attempted to parse: {content_str}")
            raise HTTPException(
                status_code=500,
                detail="Failed to parse AI response into valid JSON. Please try again."
            )
        
    except Exception as e:
        logger.error(f"Error in generate_presentation_content: {str(e)}")
        logger.error(f"Full traceback: {traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail="Failed to generate presentation content. Please try again."
        )

def create_presentation(slides_content: list) -> str:
    """Create PowerPoint presentation with professional design."""
    logger.info("Starting PowerPoint creation")
    try:
        prs = Presentation()
        
        # Set slide dimensions to widescreen (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Define common colors
        THEME_COLOR = "2F5597"  # Dark blue
        ACCENT_COLOR = "C00000"  # Dark red
        TEXT_COLOR = "000000"   # Black
        QUESTION_COLOR = "7030A0"  # Purple for questions
        ACTIVITY_COLOR = "00B050"  # Green for activities
        
        for index, slide_data in enumerate(slides_content):
            logger.info(f"Creating slide {index + 1}")
            
            if index == 0:  # Title slide
                slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add a rectangle shape as background accent
                left = 0
                top = Inches(5.5)
                width = prs.slide_width
                height = Inches(2)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(THEME_COLOR)
                shape.line.fill.background()
                
                # Add title
                title = slide.shapes.title
                title.text = slide_data["title"]
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(54)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor.from_string(TEXT_COLOR)
                title_paragraph.alignment = PP_ALIGN.CENTER
                
                # Add subtitle as a text box
                if "subtitle" in slide_data and slide_data["subtitle"]:
                    subtitle_box = slide.shapes.add_textbox(
                        Inches(1), Inches(3.5), Inches(11.333), Inches(1)
                    )
                    subtitle_frame = subtitle_box.text_frame
                    subtitle_paragraph = subtitle_frame.add_paragraph()
                    subtitle_paragraph.text = slide_data["subtitle"]
                    subtitle_paragraph.font.size = Pt(32)
                    subtitle_paragraph.font.color.rgb = RGBColor.from_string(ACCENT_COLOR)
                    subtitle_paragraph.alignment = PP_ALIGN.CENTER
            
            else:  # Content slides
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add side accent bar
                left = 0
                top = 0
                width = Inches(2)
                height = prs.slide_height
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(THEME_COLOR)
                shape.line.fill.background()
                
                # Add title box
                title_box = slide.shapes.add_textbox(
                    Inches(2.5), Inches(0.5), Inches(10), Inches(1)
                )
                title_frame = title_box.text_frame
                title_paragraph = title_frame.add_paragraph()
                title_paragraph.text = slide_data["title"]
                title_paragraph.font.size = Pt(44)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor.from_string(TEXT_COLOR)
                
                # Add subtitle if present
                top_position = Inches(1.5)
                if "subtitle" in slide_data and slide_data["subtitle"]:
                    subtitle_box = slide.shapes.add_textbox(
                        Inches(2.5), Inches(1.5), Inches(10), Inches(0.5)
                    )
                    subtitle_frame = subtitle_box.text_frame
                    subtitle_paragraph = subtitle_frame.add_paragraph()
                    subtitle_paragraph.text = slide_data["subtitle"]
                    subtitle_paragraph.font.size = Pt(28)
                    subtitle_paragraph.font.color.rgb = RGBColor.from_string(ACCENT_COLOR)
                    top_position = Inches(2)
                
                # Add content in a rounded rectangle shape
                if "content" in slide_data and slide_data["content"]:
                    content_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(2.5), top_position,
                        Inches(10), Inches(5)
                    )
                    content_shape.fill.solid()
                    content_shape.fill.fore_color.rgb = RGBColor.from_string("F2F2F2")  # Light gray
                    content_shape.line.color.rgb = RGBColor.from_string(THEME_COLOR)
                    content_shape.line.width = Pt(2)
                    
                    content_frame = content_shape.text_frame
                    content_frame.word_wrap = True
                    content_frame.margin_left = Inches(0.2)
                    content_frame.margin_right = Inches(0.2)
                    content_frame.margin_top = Inches(0.1)
                    content_frame.margin_bottom = Inches(0.1)
                    
                    # Split content into bullet points
                    bullets = slide_data["content"].split("•")
                    for bullet in bullets:
                        if bullet.strip():
                            p = content_frame.add_paragraph()
                            p.text = bullet.strip()
                            p.font.size = Pt(24)
                            
                            # Special formatting for different content types
                            if "Knowledge Check:" in bullet:
                                p.font.color.rgb = RGBColor.from_string(QUESTION_COLOR)
                                p.font.bold = True
                            elif "Think-Pair-Share" in bullet or "Activity:" in bullet:
                                p.font.color.rgb = RGBColor.from_string(ACTIVITY_COLOR)
                                p.font.bold = True
                            elif "Exit Ticket:" in bullet:
                                p.font.color.rgb = RGBColor.from_string(QUESTION_COLOR)
                                p.font.bold = True
                            else:
                                p.font.color.rgb = RGBColor.from_string(TEXT_COLOR)
                            
                            p.level = 0
                            
                # Add slide number
                slide_number = slide.shapes.add_textbox(
                    Inches(0.5), Inches(6.8), Inches(1), Inches(0.5)
                )
                slide_number_frame = slide_number.text_frame
                slide_number_paragraph = slide_number_frame.add_paragraph()
                slide_number_paragraph.text = str(index)
                slide_number_paragraph.font.size = Pt(18)
                slide_number_paragraph.font.color.rgb = RGBColor.from_string("FFFFFF")
                slide_number_paragraph.alignment = PP_ALIGN.CENTER
        
        # Save the presentation
        output_path = "generated_presentation.pptx"
        prs.save(output_path)
        logger.info("Successfully created PowerPoint presentation")
        return output_path
        
    except Exception as e:
        logger.error(f"Error in create_presentation: {str(e)}")
        logger.error(f"Full traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Error creating presentation: {str(e)}")

@app.post("/api/generate-slides")
async def generate_slides(data: Dict[str, str]):
    """Generate slides from a topic."""
    try:
        logger.info(f"Received request to generate slides for topic: {data.get('topic', 'No topic provided')}")
        
        if "topic" not in data:
            logger.error("No topic provided in request")
            raise HTTPException(status_code=400, detail="Topic is required")
        
        # Generate content using Gemini
        slides_content = await generate_presentation_content(data["topic"])
        logger.info("Content generation completed")
        
        # Create PowerPoint presentation
        output_path = create_presentation(slides_content)
        logger.info("PowerPoint creation completed")
        
        # Return the file
        logger.info("Sending response to client")
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="presentation.pptx"
        )
        
    except Exception as e:
        logger.error(f"Error in generate_slides endpoint: {str(e)}")
        logger.error(f"Full traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/health")
async def health_check():
    """Health check endpoint."""
    logger.info("Health check requested")
    api_key_configured = bool(os.getenv("GOOGLE_API_KEY"))
    logger.info(f"Google API key configured: {api_key_configured}")
    return {"status": "healthy", "google_api_key_configured": api_key_configured} 